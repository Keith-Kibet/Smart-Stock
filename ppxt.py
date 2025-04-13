from pptx import Presentation

# Create a new PowerPoint presentation
prs = Presentation()

# Function to add a slide with title and content
def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "STOCK BROKERAGE MANAGEMENT SYSTEM: A CASE OF FP-MARKETS"
subtitle.text = "Presented by: Keith Kibet\nSupervisors: Dr. Nick Ishmael\nPresentation Date: February 2025"

# Slide 2: Presentation Outline
add_slide("Presentation Outline",
          "- Introduction\n- Literature Review\n- Problem Statement\n- Justification\n"
          "- Objectives\n- Methodology\n- References\n- Thank You")

# Slide 3: Introduction
add_slide("Introduction",
          "The Stock Brokerage Management System is designed to optimize stock trading through "
          "machine learning and real-time data processing. It leverages LSTM and CNN models for "
          "stock price prediction, sentiment analysis for market trends, and risk assessment tools "
          "for better investment decisions.")

# Slide 4: Literature Review
add_slide("Literature Review",
          "Stock trading has evolved from manual to algorithmic trading and now AI-driven strategies. "
          "Machine learning models such as LSTM and CNN improve stock price forecasting and risk assessment. "
          "Sentiment analysis provides insights into market trends and trader behavior.")

# Slide 5: Problem Statement
add_slide("Problem Statement",
          "Traditional stock trading systems lack real-time data processing, efficient risk management, "
          "and adaptability to market changes. Many predictive models rely solely on historical data, "
          "resulting in delayed decision-making and increased financial risks.")

# Slide 6: Justification
add_slide("Justification",
          "The stock market is highly volatile, requiring traders to make informed decisions quickly. "
          "Traditional systems do not integrate real-time market sentiment and financial indicators effectively. "
          "This system leverages machine learning to improve accuracy, risk management, and adaptability.")

# Slide 7: Objectives
add_slide("Objectives",
          "Main Objective: Develop a Machine Learning-Driven Stock Trading Optimization System that integrates "
          "real-time stock data, predictive analytics, and risk assessment tools.\n\n"
          "Specific Objectives:\n"
          "- Provide real-time stock data for informed decisions.\n"
          "- Implement machine learning models (LSTM, CNN) for trend forecasting.\n"
          "- Incorporate risk assessment tools (Sharpe Ratio, Value at Risk).\n"
          "- Analyze sentiment data from news and social media.")

# Slide 8: Methodology
add_slide("Methodology",
          "- Data Collection: Stock data from Yahoo Finance and Finviz, sentiment data from news and social media.\n"
          "- Machine Learning Models: LSTM and CNN for stock prediction, NLP for sentiment analysis.\n"
          "- System Design: Python backend, Streamlit frontend, SQL database for authentication.\n"
          "- Deployment: Cloud-based infrastructure, Docker containers for scalability.\n"
          "- Testing: Unit testing, system testing, and security validation.")

# Slide 9: References
add_slide("References",
          "- Fischer, T., & Krauss, C. (2018). Deep learning for stock price prediction. Decision Support Systems.\n"
          "- Jiang, Z., Xu, D., & Liang, J. (2017). A deep reinforcement learning framework for portfolio management. "
          "Journal of Finance and Data Science.\n"
          "- Bollen, J., Mao, H., & Zeng, X. (2011). Twitter mood predicts the stock market. Journal of Computational Science.\n"
          "- Patel, J., Shah, S., Thakkar, P., & Kotecha, K. (2015). Predicting stock market movements using hybrid models. "
          "Expert Systems with Applications.\n"
          "- Müller, J., Schmidt, F., & Zhao, L. (2021). Algorithmic trading and machine learning approaches in financial markets. "
          "Financial Innovation Journal.")

# Slide 10: Thank You
add_slide("Thank You", "Questions?")

# Save the PowerPoint file
prs.save("D:\Personal\School Documents\y4.1\school final year project\Stock_Brokerage_Presentation.pptx")

print("Presentation slides generated successfully!")
